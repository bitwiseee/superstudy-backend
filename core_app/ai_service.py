import os
import logging
import hashlib
from django.core.cache import cache
import pymupdf as fitz  # PyMuPDF (Standardizing to use the same lib as utils.py)
from pptx import Presentation
import google.generativeai as genai
from gtts import gTTS

logger = logging.getLogger(__name__)

# --- CONFIGURATION ---
api_key = os.environ.get("GEMINI_API_KEY")
if not api_key:
    logger.warning("GEMINI_API_KEY not found in environment variables.")

# Pylance often flags 'configure' as missing due to dynamic loading. 
# The type: ignore comment suppresses this false positive.
genai.configure(api_key=api_key) # type: ignore

# Using 1.5 Flash for the large context window (1M tokens)
model = genai.GenerativeModel('gemini-1.5-flash') # type: ignore

# --- HELPERS ---

def get_stable_hash(text):
    """
    Creates a consistent hash of text for caching.
    Safely handles None or non-string inputs.
    """
    if not isinstance(text, str):
        text = ""
    return hashlib.md5(text.encode('utf-8', errors='ignore')).hexdigest()

# --- FILE EXTRACTION ---

def extract_text_from_file(file_path):
    """
    Standard text extraction for PDF, PPTX, and TXT.
    Returns stripped text or empty string on error.
    """
    if not file_path:
        return ""

    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    
    try:
        if ext == '.pdf':
            # Switched to fitz (PyMuPDF) to match core/utils.py
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += str(page.get_text()) + "\n"
        
        elif ext in ['.pptx', '.ppt']:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    # Fix: python-pptx shapes use text_frame, not .text directly
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                         text += shape.text_frame.text + "\n" # type: ignore
        
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
                
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        return ""

# --- AI FEATURES ---

def get_ai_response(document_text, user_question, language='en', document_title=""):
    """
    General Chat: Answers student questions based on the document.
    """
    # Safety Check: Ensure document_text is a string
    if document_text is None:
        document_text = ""

    text_hash = get_stable_hash(document_text[:1000])
    q_hash = get_stable_hash(user_question)
    cache_key = f"ai_chat_{text_hash}_{q_hash}_{language}"
    
    cached_response = cache.get(cache_key)
    if cached_response:
        return cached_response
    
    lang_map = {'en': 'English', 'yo': 'Yoruba', 'sw': 'Swahili', 'ha': 'Hausa', 
                'zu': 'Zulu', 'ig': 'Igbo', 'fr': 'French', 'pt': 'Portuguese'}
    lang_name = lang_map.get(language, 'English')
    
    prompt = f"""
    You are an AI tutor for African students.
    Document Title: {document_title}
    
    Context:
    {document_text[:500000]}
    
    Question: {user_question}
    
    Instructions:
    1. Answer based on the context.
    2. Be educational and encouraging.
    3. Respond in {lang_name}.
    """
    
    try:
        response = model.generate_content(prompt)
        answer = response.text
        cache.set(cache_key, answer, 3600)
        return answer
    except Exception as e:
        logger.error(f"Chat Error: {e}")
        return "Error generating response."

def generate_summary(document_text, language='en'):
    """
    Generates a concise summary of the document.
    """
    if document_text is None:
        document_text = ""

    text_hash = get_stable_hash(document_text[:1000])
    cache_key = f"ai_summary_{text_hash}_{language}"
    
    cached_summary = cache.get(cache_key)
    if cached_summary:
        return cached_summary

    lang_map = {'en': 'English', 'yo': 'Yoruba', 'sw': 'Swahili', 'ha': 'Hausa', 
                'zu': 'Zulu', 'ig': 'Igbo', 'fr': 'French', 'pt': 'Portuguese'}
    lang_name = lang_map.get(language, 'English')

    prompt = f"""
    Read the following document and provide a comprehensive summary.
    
    Document Content:
    {document_text[:500000]}
    
    Output Format:
    1. **Executive Summary**: 3-4 sentences capturing the main idea.
    2. **Key Points**: Bullet points of the most important concepts.
    3. **Conclusion**: A one-sentence takeaway.
    
    Language: {lang_name}
    """
    
    try:
        response = model.generate_content(prompt)
        summary = response.text
        cache.set(cache_key, summary, 86400) # Cache for 24 hours
        return summary
    except Exception as e:
        logger.error(f"Summary Error: {e}")
        return "Error generating summary."

def generate_flashcards(document_text, num_cards=5):
    """
    Generates flashcards (Term vs Definition) from the document.
    """
    if document_text is None:
        document_text = ""

    text_hash = get_stable_hash(document_text[:1000])
    cache_key = f"ai_flashcards_{text_hash}_{num_cards}"
    
    cached_cards = cache.get(cache_key)
    if cached_cards:
        return cached_cards

    prompt = f"""
    Create {num_cards} study flashcards based on the important concepts in this text.
    
    Document Content:
    {document_text[:500000]}
    
    Format EXACTLY as follows for parsing:
    Card 1:
    Term: [Concept Name]
    Definition: [Clear, simple explanation]
    
    Card 2:
    Term: ...
    Definition: ...
    """
    
    try:
        response = model.generate_content(prompt)
        cards = response.text
        cache.set(cache_key, cards, 86400)
        return cards
    except Exception as e:
        logger.error(f"Flashcard Error: {e}")
        return "Error generating flashcards."

def generate_quiz(document_text, num_questions=3):
    """
    Generates a multiple choice quiz based on the document.
    """
    if document_text is None:
        document_text = ""

    text_hash = get_stable_hash(document_text[:1000])
    cache_key = f"ai_quiz_{text_hash}_{num_questions}"
    
    cached_quiz = cache.get(cache_key)
    if cached_quiz:
        return cached_quiz

    prompt = f"""
    Generate a {num_questions}-question multiple choice quiz based on this text.
    
    Document Content:
    {document_text[:500000]}
    
    Format:
    Q1: [Question]
    A) [Option]
    B) [Option]
    C) [Option]
    D) [Option]
    Correct Answer: [Letter]
    Explanation: [Why it is correct]
    
    Q2: ...
    """
    
    try:
        response = model.generate_content(prompt)
        quiz = response.text
        cache.set(cache_key, quiz, 86400)
        return quiz
    except Exception as e:
        logger.error(f"Quiz Error: {e}")
        return "Error generating quiz."

# --- AUDIO GENERATION (gTTS) ---

def run_tts_sync(text, output_path, language="en"):
    """
    Converts text to speech using gTTS (Google Text-to-Speech).
    This is synchronous and will block until the file is saved.
    """
    if not isinstance(text, str) or not text.strip():
        logger.error("TTS received empty text.")
        return False

    tts_lang_map = {
        'English': 'en', 'en': 'en',
        'Yoruba': 'yo', 'yo': 'yo',
        'Hausa': 'ha', 'ha': 'ha',
        'Swahili': 'sw', 'sw': 'sw',
        'French': 'fr', 'fr': 'fr',
        'Portuguese': 'pt', 'pt': 'pt',
        # Fallbacks
        'Igbo': 'en', 'ig': 'en', 
        'Zulu': 'en', 'zu': 'en',
        'Afrikaans': 'af', 'af': 'af'
    }
    
    lang_code = tts_lang_map.get(language, 'en')
    
    try:
        tts = gTTS(text=text, lang=lang_code, slow=False)
        tts.save(output_path)
        return True
    except Exception as e:
        logger.error(f"gTTS Error: {e}")
        return False