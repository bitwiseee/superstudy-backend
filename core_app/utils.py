import fitz  # PyMuPDF
from pptx import Presentation
import os
import logging

logger = logging.getLogger(__name__)


def extract_text_from_pdf(file_path):
    """
    Extract text from PDF using PyMuPDF
    Returns formatted text with page markers
    """
    try:
        doc = fitz.open(file_path)
        text = ""
        
        # Pylance may flag doc as not iterable, but it is at runtime.
        # type: ignore suppresses the false positive.
        for page_num, page in enumerate(doc, 1): # type: ignore
            page_text = page.get_text()
            if page_text.strip():
                text += f"\n--- Page {page_num} ---\n{page_text}"
        
        doc.close()
        return text.strip()
        
    except Exception as e:
        logger.error(f"Error extracting PDF text: {e}")
        return ""


def extract_text_from_pptx(file_path):
    """
    Extract text from PowerPoint presentation
    Returns formatted text with slide markers
    """
    try:
        prs = Presentation(file_path)
        text = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"\n--- Slide {slide_num} ---\n"
            has_content = False
            
            for shape in slide.shapes:
                # FIX: Check for text_frame, not .text directly
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    # type: ignore suppresses Pylance "unknown attribute" error
                    if shape.text_frame.text.strip(): # type: ignore
                        slide_text += shape.text_frame.text + "\n" # type: ignore
                        has_content = True
            
            if has_content:
                text += slide_text
        
        return text.strip()
        
    except Exception as e:
        logger.error(f"Error extracting PPTX text: {e}")
        return ""


def extract_text_from_txt(file_path):
    """
    Extract text from plain text file
    Handles multiple encodings
    """
    encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
    
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
        except Exception as e:
            logger.error(f"Error reading text file with {encoding}: {e}")
            continue
    
    logger.error(f"Could not read text file with any encoding")
    return ""


def process_document(document):
    """
    Process uploaded document and extract text content
    
    Args:
        document: Document model instance
        
    Returns:
        str: Extracted text or empty string
    """
    file_path = document.file.path
    file_extension = os.path.splitext(file_path)[1].lower()
    
    logger.info(f"Processing document: {document.title} ({file_extension})")
    
    text = ""
    
    try:
        if file_extension == '.pdf':
            text = extract_text_from_pdf(file_path)
        elif file_extension in ['.pptx', '.ppt']:
            text = extract_text_from_pptx(file_path)
        elif file_extension == '.txt':
            text = extract_text_from_txt(file_path)
        else:
            logger.warning(f"Unsupported file format: {file_extension}")
            return ""
        
        # Clean up text
        text = clean_extracted_text(text)
        
        # Save extracted text to document
        if text:
            document.text_content = text
            document.processed = True
            document.save()
            
            logger.info(f"Successfully processed document {document.id}: {len(text)} characters")
        else:
            logger.warning(f"No text extracted from document {document.id}")
        
        return text
        
    except Exception as e:
        logger.error(f"Error processing document {document.id}: {e}")
        return ""


def clean_extracted_text(text):
    """
    Clean and normalize extracted text
    
    Args:
        text: Raw extracted text
        
    Returns:
        str: Cleaned text
    """
    if not text:
        return ""
    
    # Remove excessive whitespace
    lines = text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        # Strip whitespace
        line = line.strip()
        
        # Skip empty lines unless they're page/slide markers
        if not line:
            continue
        if line.startswith('---'):
            cleaned_lines.append(line)
        elif line:
            cleaned_lines.append(line)
    
    # Join with single newlines
    text = '\n'.join(cleaned_lines)
    
    # Remove multiple consecutive newlines (keep max 2)
    while '\n\n\n' in text:
        text = text.replace('\n\n\n', '\n\n')
    
    return text.strip()


def chunk_text(text, max_chars=8000, overlap=200):
    """
    Split text into overlapping chunks for AI processing
    
    Args:
        text: Text to chunk
        max_chars: Maximum characters per chunk
        overlap: Number of overlapping characters between chunks
        
    Returns:
        list: List of text chunks
    """
    if len(text) <= max_chars:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + max_chars
        
        # Try to break at sentence boundary
        if end < len(text):
            # Look for sentence ending
            for punct in ['. ', '.\n', '! ', '?\n']:
                last_punct = text[start:end].rfind(punct)
                if last_punct > max_chars * 0.7:  # At least 70% of chunk
                    end = start + last_punct + len(punct)
                    break
        
        chunks.append(text[start:end].strip())
        start = end - overlap if end < len(text) else end
    
    return chunks


def get_document_summary_stats(text):
    """
    Get basic statistics about document text
    
    Args:
        text: Document text
        
    Returns:
        dict: Statistics about the document
    """
    if not text:
        return {
            'word_count': 0,
            'character_count': 0,
            'page_count': 0,
            'paragraph_count': 0
        }
    
    words = text.split()
    paragraphs = [p for p in text.split('\n\n') if p.strip()]
    pages = text.count('--- Page')
    
    return {
        'word_count': len(words),
        'character_count': len(text),
        'page_count': max(1, pages),
        'paragraph_count': len(paragraphs)
    }


def validate_document_content(text, min_words=10):
    """
    Validate that extracted text is sufficient for processing
    
    Args:
        text: Extracted text
        min_words: Minimum word count required
        
    Returns:
        tuple: (is_valid, error_message)
    """
    if not text or not text.strip():
        return False, "No text could be extracted from the document"
    
    word_count = len(text.split())
    
    if word_count < min_words:
        return False, f"Document too short ({word_count} words). Minimum {min_words} words required."
    
    return True, "Valid"